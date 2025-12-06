"""
Document extractor module for Slack attachments.

Extracts text content from various document formats:
- PDF: PyPDF2
- DOCX: XML parsing (standard library, no lxml dependency) with python-docx fallback
- CSV: built-in csv module
- XLSX: openpyxl
- PPTX: XML parsing (standard library, no lxml dependency) with python-pptx fallback + LibreOffice (image conversion)
- TXT: built-in file reading

Note: DOCX and PPTX use XML parsing by default (no lxml dependency).
python-docx and python-pptx are optional fallbacks if available.
"""

import csv
import subprocess
import tempfile
import os
import json
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from io import BytesIO, StringIO
from typing import List, Optional

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError as e:
    print(json.dumps({
        "level": "WARN",
        "event": "docx_import_failed",
        "error": str(e),
        "message": "python-docx module not available - DOCX extraction will fail",
    }))
    docx = None

try:
    import openpyxl
except ImportError as e:
    print(json.dumps({
        "level": "WARN",
        "event": "openpyxl_import_failed",
        "error": str(e),
        "message": "openpyxl module not available - XLSX extraction will fail",
    }))
    openpyxl = None

try:
    from pptx import Presentation
except ImportError as e:
    print(json.dumps({
        "level": "WARN",
        "event": "pptx_import_failed",
        "error": str(e),
        "message": "python-pptx module not available - PPTX extraction will fail",
    }))
    Presentation = None


def extract_text_from_pdf(pdf_bytes: bytes) -> Optional[str]:
    """
    Extract text from PDF file.
    
    Args:
        pdf_bytes: PDF file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    if PyPDF2 is None:
        return None
    
    try:
        pdf_file = BytesIO(pdf_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_parts = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return "\n\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "pdf_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None


def _extract_text_from_docx_xml(docx_bytes: bytes) -> Optional[str]:
    """
    Extract text from DOCX file by parsing XML directly (no lxml dependency).
    
    DOCX files are ZIP archives containing XML files. This function extracts
    text from the main document XML (word/document.xml).
    
    Args:
        docx_bytes: DOCX file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    try:
        docx_zip = zipfile.ZipFile(BytesIO(docx_bytes))
        
        # DOCX namespace
        ns = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        }
        
        text_parts = []
        
        # Extract text from main document
        if 'word/document.xml' in docx_zip.namelist():
            doc_xml = docx_zip.read('word/document.xml')
            root = ET.fromstring(doc_xml)
            
            # Find all text elements (w:t)
            for text_elem in root.findall('.//w:t', ns):
                if text_elem.text:
                    text_parts.append(text_elem.text)
        
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(json.dumps({
            "level": "WARN",
            "event": "docx_xml_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None


def extract_text_from_docx(docx_bytes: bytes) -> Optional[str]:
    """
    Extract text from DOCX file.
    
    Tries python-docx first (if available), falls back to XML parsing.
    
    Args:
        docx_bytes: DOCX file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    # Try python-docx first (if available)
    if docx is not None:
        try:
            docx_file = BytesIO(docx_bytes)
            doc = docx.Document(docx_file)
            
            text_parts = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            # Extract text from headers
            for section in doc.sections:
                header = section.header
                for paragraph in header.paragraphs:
                    if paragraph.text:
                        text_parts.append(f"[Header] {paragraph.text}")
            
            # Extract text from footers
            for section in doc.sections:
                footer = section.footer
                for paragraph in footer.paragraphs:
                    if paragraph.text:
                        text_parts.append(f"[Footer] {paragraph.text}")
            
            result = "\n".join(text_parts) if text_parts else None
            if result:
                return result
        except Exception as e:
            print(json.dumps({
                "level": "WARN",
                "event": "docx_library_extraction_failed",
                "error": str(e),
                "error_type": type(e).__name__,
                "falling_back": "xml_parsing",
            }))
    
    # Fallback to XML parsing (no lxml dependency)
    return _extract_text_from_docx_xml(docx_bytes)


def extract_text_from_csv(csv_bytes: bytes) -> Optional[str]:
    """
    Extract text from CSV file.
    
    Args:
        csv_bytes: CSV file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    try:
        csv_string = csv_bytes.decode('utf-8', errors='replace')
        csv_file = StringIO(csv_string)
        reader = csv.reader(csv_file)
        
        rows = []
        for row in reader:
            rows.append(",".join(str(cell) for cell in row))
        
        return "\n".join(rows) if rows else None
    except Exception as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "csv_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None


def extract_text_from_xlsx(xlsx_bytes: bytes) -> Optional[str]:
    """
    Extract text from XLSX file.
    
    Args:
        xlsx_bytes: XLSX file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    if openpyxl is None:
        return None
    
    try:
        xlsx_file = BytesIO(xlsx_bytes)
        workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
        
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():  # Skip empty rows
                    text_parts.append(row_text)
            
            text_parts.append("")  # Separator between sheets
        
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "xlsx_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None


def _extract_text_from_pptx_xml(pptx_bytes: bytes) -> Optional[str]:
    """
    Extract text from PPTX file by parsing XML directly (no lxml dependency).
    
    PPTX files are ZIP archives containing XML files. This function extracts
    text from slide XML files (ppt/slides/slide*.xml).
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    try:
        pptx_zip = zipfile.ZipFile(BytesIO(pptx_bytes))
        
        # PPTX namespace
        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
        
        text_parts = []
        slide_num = 0
        
        # Find all slide XML files
        slide_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
        slide_files.sort()  # Process slides in order
        
        for slide_file in slide_files:
            slide_num += 1
            slide_xml = pptx_zip.read(slide_file)
            root = ET.fromstring(slide_xml)
            
            slide_texts = []
            
            # Find all text elements (a:t)
            for text_elem in root.findall('.//a:t', ns):
                if text_elem.text:
                    slide_texts.append(text_elem.text)
            
            if slide_texts:
                text_parts.append(f"Slide {slide_num}:")
                text_parts.extend(slide_texts)
                text_parts.append("")  # Separator between slides
        
        return "\n".join(text_parts) if text_parts else None
    except Exception as e:
        print(json.dumps({
            "level": "WARN",
            "event": "pptx_xml_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None


def extract_text_from_pptx(pptx_bytes: bytes) -> Optional[str]:
    """
    Extract text from PPTX file.
    
    Tries python-pptx first (if available), falls back to XML parsing.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    # Try python-pptx first (if available)
    if Presentation is not None:
        try:
            pptx_file = BytesIO(pptx_bytes)
            prs = Presentation(pptx_file)
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Text boxes and placeholders
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                    
                    # Tables
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_texts.append(" | ".join(row_text))
                    
                    # Group shapes (nested shapes)
                    if hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text:
                                slide_texts.append(sub_shape.text)
                
                if slide_texts:
                    text_parts.append(f"Slide {slide_num}:")
                    text_parts.extend(slide_texts)
                    text_parts.append("")  # Separator between slides
            
            result = "\n".join(text_parts) if text_parts else None
            if result:
                return result
        except Exception as e:
            print(json.dumps({
                "level": "WARN",
                "event": "pptx_library_extraction_failed",
                "error": str(e),
                "error_type": type(e).__name__,
                "falling_back": "xml_parsing",
            }))
    
    # Fallback to XML parsing (no lxml dependency)
    return _extract_text_from_pptx_xml(pptx_bytes)


def convert_pptx_slides_to_images(pptx_bytes: bytes) -> Optional[List[bytes]]:
    """
    Convert PPTX slides to PNG images using LibreOffice.
    
    Requires LibreOffice installed in Lambda Layer at /opt/libreoffice.
    Implements error handling, timeout (60 seconds), and automatic cleanup.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        List of image bytes (one per slide), or None if conversion fails
    """
    # Save PPTX to temporary file
    tmp_pptx_path = None
    output_dir = None
    
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
            tmp_pptx.write(pptx_bytes)
            tmp_pptx_path = tmp_pptx.name
        
        # Use LibreOffice to convert PPTX to images
        output_dir = tempfile.mkdtemp()
        
        # LibreOffice path in Lambda Layer
        libreoffice_path = "/opt/libreoffice/program/soffice"
        if not os.path.exists(libreoffice_path):
            # Fallback to system LibreOffice (for local testing)
            libreoffice_path = "libreoffice"
            print(json.dumps({
                "level": "WARN",
                "event": "libreoffice_fallback_to_system",
                "message": "LibreOffice Lambda Layer not found, using system LibreOffice",
            }))
        
        # Convert PPTX to PNG images (one per slide) with 60 second timeout
        print(json.dumps({
            "level": "INFO",
            "event": "pptx_conversion_started",
            "pptx_size": len(pptx_bytes),
            "timeout_seconds": 60,
        }))
        
        result = subprocess.run([
            libreoffice_path,
            '--headless',
            '--convert-to', 'png',
            '--outdir', output_dir,
            tmp_pptx_path
        ], check=True, timeout=60, capture_output=True, text=True)
        
        # Read generated images
        image_files = sorted(Path(output_dir).glob('*.png'))
        images = []
        for img_file in image_files:
            with open(img_file, 'rb') as f:
                images.append(f.read())
        
        print(json.dumps({
            "level": "INFO",
            "event": "pptx_conversion_success",
            "slide_count": len(images),
        }))
        
        return images if images else None
    except subprocess.TimeoutExpired:
        print(json.dumps({
            "level": "ERROR",
            "event": "pptx_conversion_timeout",
            "timeout_seconds": 60,
            "message": "LibreOffice conversion exceeded 60 second timeout",
        }))
        return None
    except subprocess.CalledProcessError as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "pptx_conversion_failed",
            "error": str(e),
            "returncode": e.returncode,
            "stdout": e.stdout[:500] if e.stdout else None,
            "stderr": e.stderr[:500] if e.stderr else None,
            "message": "LibreOffice conversion process failed",
        }))
        return None
    except Exception as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "pptx_conversion_unexpected_error",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None
    finally:
        # Cleanup temporary files (always executed)
        if tmp_pptx_path and os.path.exists(tmp_pptx_path):
            try:
                os.unlink(tmp_pptx_path)
            except Exception as cleanup_error:
                print(json.dumps({
                    "level": "WARN",
                    "event": "pptx_cleanup_failed",
                    "file": tmp_pptx_path,
                    "error": str(cleanup_error),
                }))
        if output_dir and os.path.exists(output_dir):
            try:
                shutil.rmtree(output_dir, ignore_errors=True)
            except Exception as cleanup_error:
                print(json.dumps({
                    "level": "WARN",
                    "event": "pptx_output_cleanup_failed",
                    "directory": output_dir,
                    "error": str(cleanup_error),
                }))


def extract_text_from_txt(txt_bytes: bytes) -> Optional[str]:
    """
    Extract text from TXT file.
    
    Args:
        txt_bytes: TXT file content as bytes
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    try:
        return txt_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        print(json.dumps({
            "level": "ERROR",
            "event": "txt_extraction_failed",
            "error": str(e),
            "error_type": type(e).__name__,
        }))
        return None

