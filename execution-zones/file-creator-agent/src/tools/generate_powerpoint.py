"""
generate_powerpoint tool for Execution Agent (027-slack-file-generation-best-practices).

Produces PowerPoint (.pptx) presentations and stores them in tool_context.invocation_state
for the handler to extract and build file_artifact.
"""

from io import BytesIO
from typing import Any, List

from strands import tool

import file_config as fc
from file_config import sanitize_filename

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@tool(context=True)
def generate_powerpoint(
    filename: str,
    slides: List[dict],
    tool_context: Any,
) -> str:
    """PowerPointプレゼンテーション (.pptx) を生成します。

    スライドごとにタイトルと本文を指定できます。

    Args:
        filename: ファイル名（拡張子不要、自動付加）
        slides: スライドの配列。各要素は title, body を持つ。layout はオプション。
        tool_context: Strands framework context (injected). Contains invocation_state.

    Returns:
        モデルに返す説明文（日本語）。ファイルは invocation_state に格納される。
    """
    if not filename or not isinstance(filename, str) or not filename.strip():
        return "エラー: filename を指定してください。"
    if not slides or not isinstance(slides, list):
        return "エラー: slides は必須です（title, body を持つオブジェクトの配列）。"

    if len(slides) > 200:
        return "エラー: スライド数は最大200までです。"

    safe_filename = sanitize_filename(filename.strip(), "pptx")
    if not safe_filename or "." not in safe_filename:
        safe_filename = f"{safe_filename}.pptx"

    try:
        from pptx import Presentation
    except ImportError:
        return "エラー: python-pptx がインストールされていません。"

    prs = Presentation()

    layout_map = {
        "title_slide": 0,
        "title_and_content": 1,
        "blank": 6,
    }

    for slide_def in slides:
        if not isinstance(slide_def, dict):
            continue
        title = slide_def.get("title") or ""
        body = slide_def.get("body") or ""
        layout_name = slide_def.get("layout") or "title_and_content"
        layout_idx = layout_map.get(layout_name, 1)

        slide_layouts = prs.slide_layouts
        if layout_idx >= len(slide_layouts):
            layout_idx = 1
        slide = prs.slides.add_slide(slide_layouts[layout_idx])

        if layout_idx == 6:
            continue
        if slide.shapes.title:
            slide.shapes.title.text = str(title)
        if layout_idx == 1 and len(slide.shapes.placeholders) > 1:
            body_placeholder = slide.shapes.placeholders[1]
            if body_placeholder.has_text_frame:
                body_placeholder.text_frame.text = str(body)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    file_bytes = buf.getvalue()

    if len(file_bytes) > fc.MAX_OFFICE_FILE_BYTES:
        return (
            f"エラー: ファイルサイズが上限（{fc.MAX_OFFICE_FILE_BYTES} バイト）を超えています。"
            f"現在 {len(file_bytes)} バイトです。"
        )

    invocation_state = getattr(tool_context, "invocation_state", {})
    if isinstance(invocation_state, dict):
        invocation_state["generated_file"] = {
            "file_bytes": file_bytes,
            "file_name": safe_filename,
            "mime_type": _PPTX_MIME,
            "description": f"PowerPointプレゼンテーション「{safe_filename}」を生成しました。",
        }

    return f"ファイル「{safe_filename}」を作成しました。Slack にアップロードされます。"
